"""Tools del Slides Agent: presentaciones .pptx desde contenido de Looker.

Enfoque: las slides son narrativa visual, no tablas. Cada tile del dashboard
se renderiza como PNG (query render task de Looker) y ocupa una lámina con su
título. El template corporativo (.pptx en el bucket) aporta el master, la
tipografía y la lámina de portada.
"""
import io
import json

from pptx import Presentation
from pptx.util import Emu, Inches, Pt

from common.delivery import SIGNED_URL_HOURS, timestamped_name, upload_and_sign
from common.looker_client import get_sdk, poll_render_task

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)  # 16:9
IMG_W, IMG_H = 1400, 760                        # px del render por tile


def _base_presentation(template_name: str = "") -> Presentation:
    if not template_name:
        prs = Presentation()
        prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
        return prs
    from common import templates
    name = template_name if template_name.endswith(".pptx") else f"{template_name}.pptx"
    return Presentation(io.BytesIO(templates.load_bytes("slides", name)))


def _blank_layout(prs: Presentation):
    # el layout con menos placeholders disponible
    return min(prs.slide_layouts, key=lambda l: len(l.placeholders))


def _add_title_slide(prs: Presentation, title: str, subtitle: str):
    slide = prs.slides.add_slide(_blank_layout(prs))
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.4), prs.slide_width - Inches(1.6), Inches(1.6))
    tf = box.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(40)
    tf.paragraphs[0].font.bold = True
    p = tf.add_paragraph()
    p.text = subtitle
    p.font.size = Pt(16)


def _add_image_slide(prs: Presentation, title: str, png: bytes):
    slide = prs.slides.add_slide(_blank_layout(prs))
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), prs.slide_width - Inches(1.0), Inches(0.7))
    tf = box.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True
    pic = slide.shapes.add_picture(io.BytesIO(png), 0, Inches(1.1),
                                   width=prs.slide_width - Inches(1.0))
    pic.left = int((prs.slide_width - pic.width) / 2)
    max_h = prs.slide_height - Inches(1.4)
    if pic.height > max_h:  # reescalar si desborda
        ratio = max_h / pic.height
        pic.height, pic.width = int(pic.height * ratio), int(pic.width * ratio)
        pic.left = int((prs.slide_width - pic.width) / 2)
        pic.top = Inches(1.1)


def _render_query_png(sdk, query_id: str) -> bytes:
    task = sdk.create_query_render_task(query_id=query_id, result_format="png",
                                        width=IMG_W, height=IMG_H)
    return poll_render_task(sdk, task.id)


def export_dashboard_to_slides(dashboard_id: str, filename: str = "",
                               template_name: str = "", subtitle: str = "") -> str:
    """Genera una presentación .pptx desde un dashboard: portada + una lámina
    por tile (título del tile + su visualización renderizada como imagen).
    template_name: .pptx corporativo del bucket de templates (opcional)."""
    sdk = get_sdk()
    db = sdk.dashboard(dashboard_id=dashboard_id, fields="id,title,dashboard_elements")
    prs = _base_presentation(template_name)
    _add_title_slide(prs, db.title or "Dashboard",
                     subtitle or f"Generado desde Looker · dashboard {dashboard_id}")
    exported, skipped = [], []
    for el in (db.dashboard_elements or []):
        if el.type != "vis" or not el.query_id:
            continue
        try:
            png = _render_query_png(sdk, el.query_id)
            _add_image_slide(prs, el.title or f"Tile {el.id}", png)
            exported.append(el.title or el.id)
        except Exception as e:  # un tile que no renderiza no tumba la presentación
            skipped.append({"tile": el.title or el.id, "error": str(e)})
    if not exported:
        return json.dumps({"error": "Ningún tile del dashboard pudo renderizarse.",
                           "skipped": skipped}, ensure_ascii=False)
    buf = io.BytesIO()
    prs.save(buf)
    url = upload_and_sign(buf.getvalue(),
                          timestamped_name(filename or db.title or "presentacion", ".pptx"))
    return json.dumps({"download_url": url, "slides": len(exported) + 1,
                       "tiles": exported, "skipped": skipped,
                       "expires_in_hours": SIGNED_URL_HOURS}, ensure_ascii=False)


def list_slides_templates() -> str:
    """Lista los templates .pptx corporativos disponibles."""
    from common import templates
    return json.dumps(templates.list_names("slides"))


ALL_TOOLS = [export_dashboard_to_slides, list_slides_templates]
